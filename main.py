import os
import tempfile
import logging
import json
from datetime import datetime
from flask import Flask, request, jsonify
import firebase_admin
from firebase_admin import credentials, db
from google.cloud import storage
from google.cloud import aiplatform
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_vertexai import VertexAIEmbeddings
import PyPDF2
from PIL import Image
import pytesseract
import docx
import pptx
from pydub import AudioSegment
import speech_recognition as sr

app = Flask(__name__)
logging.basicConfig(level=logging.INFO)

# --- Configuration ---
PROJECT_ID = os.environ.get("PROJECT_ID")
GCP_REGION = os.environ.get("GCP_REGION")
EMBEDDING_MODEL = os.environ.get("EMBEDDING_MODEL")
INDEX_ID = os.environ.get("INDEX_ID")
STAGING_BUCKET_URI = os.environ.get("STAGING_BUCKET_URI")

# Check for required environment variables
if not all([PROJECT_ID, GCP_REGION, EMBEDDING_MODEL, INDEX_ID, STAGING_BUCKET_URI]):
    raise ValueError("Missing one or more required environment variables")

# Initialize Firebase
if not firebase_admin._apps:
    cred = credentials.ApplicationDefault()
    firebase_admin.initialize_app(cred, {
        'databaseURL': f"https://{PROJECT_ID}-default-rtdb.firebaseio.com"
    })

# Initialize services
storage_client = storage.Client()
aiplatform.init(project=PROJECT_ID, location=GCP_REGION, staging_bucket=STAGING_BUCKET_URI)

embeddings = VertexAIEmbeddings(
    model_name=EMBEDDING_MODEL,
    project=PROJECT_ID,
    location=GCP_REGION
)

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000,
    chunk_overlap=200
)

class MatchingEngineHelper:
    def __init__(self, index_id: str):
        self.index = aiplatform.MatchingEngineIndex(index_name=index_id)

    def add_embeddings(self, texts: list[str], metadatas: list[dict]) -> bool:
        try:
            logging.info(f"Generating embeddings for {len(texts)} texts...")
            embeddings_list = embeddings.embed_documents(texts)
            logging.info("Embeddings generated successfully.")

            bucket_name = STAGING_BUCKET_URI.replace("gs://", "")
            bucket = storage_client.bucket(bucket_name)
            json_objects = []

            for i, text in enumerate(texts):
                metadata = metadatas[i]
                doc_id = metadata['doc_id']
                chunk_index = metadata['chunk_index']
                
                # Upload chunk to GCS in the staging bucket
                chunk_filename = f"chunks/{doc_id}_{chunk_index}.txt"
                logging.info(f"Uploading chunk to GCS: gs://{bucket_name}/{chunk_filename}")
                blob = bucket.blob(chunk_filename)
                try:
                    blob.upload_from_string(text, content_type="text/plain")
                    logging.info(f"Successfully uploaded chunk: {chunk_filename}")
                except Exception as e:
                    logging.error(f"Failed to upload chunk {chunk_filename}: {str(e)}")
                    raise e
                
                # Create a dictionary of metadata to be stored with the embedding
                embedding_metadata = {
                    "doc_id": doc_id,
                    "chunk_index": chunk_index,
                    "deal_id": metadata['deal_id'],
                    "text": text,
                    "filename": metadata['filename'],
                    "file_ext": metadata['file_ext']
                }

                json_objects.append({
                    "id": f"{doc_id}_{chunk_index}",
                    "embedding": embeddings_list[i],
                    "restricts": [
                        {"namespace": "deal_id", "allow": [metadata['deal_id']]}
                    ],
                    "numeric_restricts": [
                        {"namespace": "chunk_index", "value_int": chunk_index}
                    ],
                    "crowding_tag": {"crowding_attribute": doc_id},
                    "metadata": embedding_metadata
                })

            # Create a JSONL file in a temporary directory for batch update
            with tempfile.NamedTemporaryFile(mode='w+', delete=False, suffix=".json") as temp_json_file:
                for obj in json_objects:
                    temp_json_file.write(json.dumps(obj) + "\n")
                temp_file_path = temp_json_file.name
        
            # Upload the JSONL file to GCS
            blob_name = f"embeddings/{os.path.basename(temp_file_path)}"
            blob = bucket.blob(blob_name)
            blob.upload_from_filename(temp_file_path)
            gcs_uri = f"gs://{bucket_name}/{blob_name}"
            logging.info(f"Uploaded embeddings data to {gcs_uri}")

            # Update the index
            self.index.update_embeddings(
                contents_delta_uri=os.path.dirname(gcs_uri),
            )
            logging.info(f"Successfully triggered index update for {self.index.name}")
            
            return True
        except Exception as e:
            logging.error(f"Error adding embeddings to Vertex AI Vector Search: {str(e)}")
            return False
        finally:
            if 'temp_file_path' in locals() and os.path.exists(temp_file_path):
                os.unlink(temp_file_path)

matching_engine_helper = MatchingEngineHelper(index_id=INDEX_ID)

class DocumentProcessor:
    def _extract_text_from_pdf(self, file_path):
        text = ""
        metadata = {}
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() + '\n'
            doc_info = reader.metadata
            metadata = {
                "author": doc_info.author,
                "creator": doc_info.creator,
                "producer": doc_info.producer,
                "subject": doc_info.subject,
                "title": doc_info.title,
                "pages": len(reader.pages)
            }
        return text, metadata

    def _extract_text_from_docx(self, file_path):
        doc = docx.Document(file_path)
        text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        props = doc.core_properties
        metadata = {
            "author": props.author,
            "created": props.created.isoformat() if props.created else None,
            "last_modified_by": props.last_modified_by,
            "last_printed": props.last_printed.isoformat() if props.last_printed else None,
            "modified": props.modified.isoformat() if props.modified else None,
            "title": props.title,
            "version": props.version
        }
        return text, metadata

    def _extract_text_from_pptx(self, file_path):
        prs = pptx.Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        props = prs.core_properties
        metadata = {
            "author": props.author,
            "created": props.created.isoformat() if props.created else None,
            "last_modified_by": props.last_modified_by,
            "modified": props.modified.isoformat() if props.modified else None,
            "title": props.title,
            "version": props.version
        }
        return text, metadata

    def _extract_text_from_image(self, file_path):
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        return text, {}

    def _extract_text_from_audio(self, file_path, file_ext):
        if file_ext in ['.mp4', '.mov', '.avi']:
            audio = AudioSegment.from_file(file_path)
            audio_path = tempfile.mktemp(suffix=".wav")
            audio.export(audio_path, format="wav")
        else:
            audio_path = file_path
        
        r = sr.Recognizer()
        with sr.AudioFile(audio_path) as source:
            audio_data = r.record(source)
            text = r.recognize_google(audio_data)
        
        if audio_path != file_path:
            os.unlink(audio_path)
            
        return text, {"transcription_engine": "google_web_speech_api"}

    def extract_text_and_metadata(self, file_path, filename):
        """Extract text and metadata from various file types"""
        file_ext = os.path.splitext(filename)[1].lower()
        text, metadata = "", {}
        
        try:
            if file_ext == '.pdf':
                text, metadata = self._extract_text_from_pdf(file_path)
            elif file_ext == '.txt':
                with open(file_path, 'r', encoding='utf-8') as file:
                    text = file.read()
            elif file_ext in ['.png', '.jpg', '.jpeg']:
                text, metadata = self._extract_text_from_image(file_path)
            elif file_ext == '.docx':
                text, metadata = self._extract_text_from_docx(file_path)
            elif file_ext == '.pptx':
                text, metadata = self._extract_text_from_pptx(file_path)
            elif file_ext in ['.mp3', '.wav', '.m4a', '.mp4', '.mov', '.avi']:
                text, metadata = self._extract_text_from_audio(file_path, file_ext)
            else:
                logging.warning(f"Unsupported file type: {file_ext}")
                text = "Unsupported file type"
        except Exception as e:
            logging.error(f"Error extracting text from {filename}: {str(e)}")
            text = ""

        # Add filename to metadata
        metadata['filename'] = filename
        metadata['file_ext'] = file_ext
        return text, metadata

    def process_document(self, doc_id, bucket_name, filename, deal_id):
        """Main processing pipeline"""
        temp_file_path = None
        try:
            logging.info(f"Starting processing for document: {doc_id}")
            db.reference(f'documents/{doc_id}').update({'status': 'processing'})

            logging.info(f"Downloading document: gs://{bucket_name}/{filename}")
            bucket = storage_client.bucket(bucket_name)
            blob = bucket.blob(filename)

            with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(filename)[1]) as temp_file:
                temp_file_path = temp_file.name
                blob.download_to_filename(temp_file_path)
            logging.info(f"Successfully downloaded document to: {temp_file_path}")

            logging.info(f"Extracting text and metadata from: {filename}")
            text_content, metadata = self.extract_text_and_metadata(temp_file_path, filename)
            logging.info(f"Successfully extracted text and metadata.")

            metadata['deal_id'] = deal_id

            if not text_content or not text_content.strip():
                logging.warning(f"No text extracted from {filename}. Marking as failed.")
                db.reference(f'documents/{doc_id}').update({
                    'status': 'failed',
                    'error': 'No text extracted from document',
                    'failed_at': datetime.utcnow().isoformat()
                })
                return False

            logging.info(f"Splitting text into chunks for document: {doc_id}")
            chunks = text_splitter.split_text(text_content)
            logging.info(f"Created {len(chunks)} chunks for document {doc_id}")

            metadatas_to_add = []
            for i, chunk in enumerate(chunks):
                chunk_metadata = metadata.copy()
                chunk_metadata.update({
                    'doc_id': doc_id,
                    'chunk_index': i,
                    'created_at': datetime.utcnow().isoformat()
                })
                metadatas_to_add.append(chunk_metadata)

            logging.info(f"Adding embeddings to Vector Search for document: {doc_id}")
            embeddings_stored_status = matching_engine_helper.add_embeddings(
                texts=chunks,
                metadatas=metadatas_to_add
            )

            if embeddings_stored_status:
                logging.info(f"Successfully added embeddings for document: {doc_id}")
                db.reference(f'documents/{doc_id}').update({
                    'status': 'completed',
                    'processed_content': text_content[:500] + ('...' if len(text_content) > 500 else ''),
                    'embeddings_stored': True,
                    'chunk_count': len(chunks),
                    'processed_at': datetime.utcnow().isoformat(),
                    'metadata': metadata
                })
            else:
                logging.error(f"Failed to add embeddings for document: {doc_id}")
                db.reference(f'documents/{doc_id}').update({
                    'status': 'failed',
                    'error': 'Failed to add embeddings to Vector Search',
                    'failed_at': datetime.utcnow().isoformat()
                })

            logging.info(f"Successfully processed document {doc_id}")
            return True

        except Exception as e:
            logging.error(f"An unexpected error occurred during processing for doc_id: {doc_id}: {str(e)}")
            db.reference(f'documents/{doc_id}').update({
                'status': 'failed',
                'error': str(e),
                'failed_at': datetime.utcnow().isoformat()
            })
            return False
        finally:
            if temp_file_path and os.path.exists(temp_file_path):
                os.unlink(temp_file_path)

processor = DocumentProcessor()

@app.route('/health', methods=['GET'])
def health_check():
    return jsonify({'status': 'healthy'})

@app.route('/process', methods=['POST'])
def process_document_endpoint():
    """Process document endpoint called by Cloud Function"""
    data = request.get_json()
    
    doc_id = data.get('doc_id')
    bucket_name = data.get('bucket')
    filename = data.get('filename')
    deal_id = data.get('deal_id')
    
    if not all([doc_id, bucket_name, filename, deal_id]):
        return jsonify({'error': 'Missing required parameters'}), 400
    
    logging.info(f"Processing request for doc_id: {doc_id}, file: {filename}, deal_id: {deal_id}")
    
    success = processor.process_document(doc_id, bucket_name, filename, deal_id)
    
    if success:
        return jsonify({'status': 'success', 'doc_id': doc_id})
    else:
        return jsonify({'status': 'failed', 'doc_id': doc_id}), 500

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 8080))
    app.run(host='0.0.0.0', port=port, debug=True)